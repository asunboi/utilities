from __future__ import annotations

import argparse
import ast
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

EMU_PER_INCH = 914400
TRACKER_SHAPE_PREFIX = "DeckDots Tracker"
ACTIVE_COLOR = RGBColor(0x00, 0x00, 0x00)
INACTIVE_COLOR = RGBColor(0xE7, 0xE6, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LEFT_MARGIN = int(0.315 * EMU_PER_INCH)
RIGHT_MARGIN = int(0.35 * EMU_PER_INCH)
BOTTOM_MARGIN = int(0.08 * EMU_PER_INCH)
TOPIC_GAP = int(0.60 * EMU_PER_INCH)
LABEL_TO_DOTS_GAP = int(0.05 * EMU_PER_INCH)
LABEL_BOX_HEIGHT = int(0.22 * EMU_PER_INCH)
LABEL_FONT_SIZE = Pt(12)
TEXT_BOX_MARGIN_X = int(0.04 * EMU_PER_INCH)
TEXT_BOX_MARGIN_Y = int(0.02 * EMU_PER_INCH)


@dataclass(frozen=True)
class TopicConfig:
    name: str
    slides_per_dot: tuple[int, ...]

    @property
    def dot_count(self) -> int:
        return len(self.slides_per_dot)

    @property
    def slide_count(self) -> int:
        return sum(self.slides_per_dot)


@dataclass(frozen=True)
class SlideState:
    topic_index: int
    dot_index: int


@dataclass(frozen=True)
class TopicLayout:
    x: int
    width: int
    dot_positions: tuple[int, ...]
    dot_diameter: int
    label_y: int
    dots_y: int


def default_output_path(input_path: Path) -> Path:
    return input_path.with_name(f"{input_path.stem}.progress-dots{input_path.suffix}")


def parse_scalar(value: str) -> str:
    value = value.strip()
    if not value:
        raise ValueError("Expected a value in config.")
    if value[0] in {'"', "'"}:
        return ast.literal_eval(value)
    return value


def parse_int_list(value: str) -> list[int]:
    parsed = ast.literal_eval(value)
    if not isinstance(parsed, list):
        raise ValueError("Expected a list of integers for slides_per_dot.")
    return [int(item) for item in parsed]


def parse_simple_yaml_config(text: str) -> dict:
    topics: list[dict] = []
    current_topic: dict | None = None
    collecting_multiline_dots = False

    for raw_line in text.splitlines():
        line = raw_line.split("#", 1)[0].rstrip()
        if not line.strip():
            continue

        stripped = line.strip()
        indent = len(line) - len(line.lstrip(" "))

        if stripped == "topics:":
            continue

        if indent == 2 and stripped.startswith("- name:"):
            current_topic = {"name": parse_scalar(stripped[len("- name:") :]), "slides_per_dot": []}
            topics.append(current_topic)
            collecting_multiline_dots = False
            continue

        if current_topic is None:
            raise ValueError("Config must define topics before topic fields.")

        if indent == 4 and stripped.startswith("slides_per_dot:"):
            remainder = stripped[len("slides_per_dot:") :].strip()
            if remainder:
                current_topic["slides_per_dot"] = parse_int_list(remainder)
                collecting_multiline_dots = False
            else:
                current_topic["slides_per_dot"] = []
                collecting_multiline_dots = True
            continue

        if collecting_multiline_dots and indent == 6 and stripped.startswith("- "):
            current_topic["slides_per_dot"].append(int(stripped[2:].strip()))
            continue

        raise ValueError(f"Unsupported config line: {raw_line}")

    return {"topics": topics}


def load_raw_config(config_path: Path) -> dict:
    text = config_path.read_text(encoding="utf-8")
    try:
        import yaml  # type: ignore
    except ModuleNotFoundError:
        return parse_simple_yaml_config(text)

    raw = yaml.safe_load(text)
    if not isinstance(raw, dict):
        raise ValueError("Config root must be a mapping.")
    return raw


def load_topics(config_path: Path) -> list[TopicConfig]:
    raw_config = load_raw_config(config_path)
    raw_topics = raw_config.get("topics")
    if not isinstance(raw_topics, list) or not raw_topics:
        raise ValueError("Config must include a non-empty 'topics' list.")

    topics: list[TopicConfig] = []
    for raw_topic in raw_topics:
        if not isinstance(raw_topic, dict):
            raise ValueError("Each topic must be a mapping.")
        name = raw_topic.get("name")
        slides_per_dot = raw_topic.get("slides_per_dot")
        if not isinstance(name, str) or not name.strip():
            raise ValueError("Each topic must include a non-empty name.")
        if not isinstance(slides_per_dot, list) or not slides_per_dot:
            raise ValueError(f"Topic '{name}' must include a non-empty slides_per_dot list.")

        normalized_counts = tuple(int(count) for count in slides_per_dot)
        if any(count <= 0 for count in normalized_counts):
            raise ValueError(f"Topic '{name}' has a non-positive slide count.")
        topics.append(TopicConfig(name=name.strip(), slides_per_dot=normalized_counts))

    return topics


def build_slide_states(topics: list[TopicConfig]) -> list[SlideState]:
    states: list[SlideState] = []
    for topic_index, topic in enumerate(topics):
        for dot_index, slide_span in enumerate(topic.slides_per_dot):
            states.extend(
                SlideState(topic_index=topic_index, dot_index=dot_index)
                for _ in range(slide_span)
            )
    return states


def estimate_label_width(name: str) -> int:
    average_char_width = 0.095 * EMU_PER_INCH
    return int(len(name) * average_char_width) + 2 * TEXT_BOX_MARGIN_X


def fit_grouped_layout(
    slide_width: int,
    slide_height: int,
    topics: list[TopicConfig],
) -> list[TopicLayout]:
    available_width = slide_width - LEFT_MARGIN - RIGHT_MARGIN
    topic_gap = TOPIC_GAP
    dot_diameter = int(0.115 * EMU_PER_INCH)
    dot_gap = int(0.05 * EMU_PER_INCH)

    while True:
        layouts: list[TopicLayout] = []
        cursor_x = LEFT_MARGIN
        label_y = slide_height - BOTTOM_MARGIN - dot_diameter - LABEL_TO_DOTS_GAP - LABEL_BOX_HEIGHT
        dots_y = slide_height - BOTTOM_MARGIN - dot_diameter

        for topic in topics:
            dot_positions = tuple(
                cursor_x + dot_index * (dot_diameter + dot_gap)
                for dot_index in range(topic.dot_count)
            )
            dots_width = topic.dot_count * dot_diameter + max(topic.dot_count - 1, 0) * dot_gap
            label_width = max(estimate_label_width(topic.name), dots_width)
            group_width = label_width
            layouts.append(
                TopicLayout(
                    x=cursor_x,
                    width=group_width,
                    dot_positions=dot_positions,
                    dot_diameter=dot_diameter,
                    label_y=label_y,
                    dots_y=dots_y,
                )
            )
            cursor_x += group_width + topic_gap

        used_width = cursor_x - LEFT_MARGIN - topic_gap
        if used_width <= available_width or (dot_diameter <= int(0.08 * EMU_PER_INCH) and topic_gap <= int(0.34 * EMU_PER_INCH)):
            return layouts

        dot_diameter = max(int(dot_diameter * 0.94), int(0.08 * EMU_PER_INCH))
        dot_gap = max(int(dot_gap * 0.92), int(0.03 * EMU_PER_INCH))
        topic_gap = max(int(topic_gap * 0.92), int(0.34 * EMU_PER_INCH))


def is_tracker_shape(shape) -> bool:
    return getattr(shape, "name", "").startswith(TRACKER_SHAPE_PREFIX)


def remove_existing_tracker_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if not is_tracker_shape(shape):
            continue
        element = shape._element
        element.getparent().remove(element)


def add_topic_label(slide, topic_name: str, layout: TopicLayout, active: bool) -> None:
    label_box = slide.shapes.add_textbox(layout.x, layout.label_y, layout.width, LABEL_BOX_HEIGHT)
    label_box.name = f"{TRACKER_SHAPE_PREFIX} Topic Label {topic_name}"
    text_frame = label_box.text_frame
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = TEXT_BOX_MARGIN_X
    text_frame.margin_right = TEXT_BOX_MARGIN_X
    text_frame.margin_top = TEXT_BOX_MARGIN_Y
    text_frame.margin_bottom = TEXT_BOX_MARGIN_Y

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = topic_name
    run.font.name = "Arial"
    run.font.size = LABEL_FONT_SIZE
    run.font.color.rgb = ACTIVE_COLOR if active else INACTIVE_COLOR


def add_topic_dots(
    slide,
    topic_name: str,
    layout: TopicLayout,
    dot_count: int,
    active_topic: bool,
    active_dot_index: int | None,
) -> None:
    for dot_index in range(dot_count):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            layout.dot_positions[dot_index],
            layout.dots_y,
            layout.dot_diameter,
            layout.dot_diameter,
        )
        dot.name = f"{TRACKER_SHAPE_PREFIX} Topic Dot {topic_name} {dot_index + 1}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACTIVE_COLOR if active_topic and dot_index == active_dot_index else WHITE
        dot.line.color.rgb = ACTIVE_COLOR if active_topic else INACTIVE_COLOR
        dot.line.width = Pt(1)


def add_progress_footer(
    slide,
    slide_state: SlideState,
    topics: list[TopicConfig],
    slide_width: int,
    slide_height: int,
) -> None:
    remove_existing_tracker_shapes(slide)
    layouts = fit_grouped_layout(slide_width, slide_height, topics)

    for topic_index, topic in enumerate(topics):
        layout = layouts[topic_index]
        active_topic = topic_index == slide_state.topic_index
        add_topic_label(slide, topic.name, layout, active_topic)
        add_topic_dots(
            slide=slide,
            topic_name=topic.name,
            layout=layout,
            dot_count=topic.dot_count,
            active_topic=active_topic,
            active_dot_index=slide_state.dot_index if active_topic else None,
        )


def add_progress_dots_to_presentation(input_path: Path, config_path: Path, output_path: Path) -> int:
    topics = load_topics(config_path)
    presentation = Presentation(str(input_path))
    slide_states = build_slide_states(topics)

    if len(slide_states) != len(presentation.slides):
        raise ValueError(
            f"Config describes {len(slide_states)} slides but presentation has {len(presentation.slides)} slides."
        )

    for slide_index, slide in enumerate(presentation.slides):
        add_progress_footer(
            slide=slide,
            slide_state=slide_states[slide_index],
            topics=topics,
            slide_width=presentation.slide_width,
            slide_height=presentation.slide_height,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return len(presentation.slides)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Add a grouped bottom-row topic tracker to every slide in a PowerPoint deck."
    )
    parser.add_argument("input_pptx", type=Path, help="Path to the input .pptx file.")
    parser.add_argument(
        "--config",
        required=True,
        type=Path,
        help="Path to the tracker config .yaml file.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        help="Optional output path. Defaults to '<input>.progress-dots.pptx'.",
    )
    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()

    input_path = args.input_pptx.resolve()
    config_path = args.config.resolve()
    output_path = (args.output or default_output_path(input_path)).resolve()

    slide_count = add_progress_dots_to_presentation(input_path, config_path, output_path)
    print(f"Wrote {output_path} with grouped progress dots on {slide_count} slides.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
