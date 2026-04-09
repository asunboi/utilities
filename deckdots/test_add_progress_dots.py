from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from deckdots.add_progress_dots import (
    TRACKER_SHAPE_PREFIX,
    add_progress_dots_to_presentation,
    build_slide_states,
    load_topics,
)

SAMPLE_DECK = Path(__file__).with_name("260406_okr.pptx")
SAMPLE_CONFIG = Path(__file__).with_name("configs") / "260406_okr.yaml"


class AddProgressDotsTest(unittest.TestCase):
    def test_loads_topic_config_and_expands_repeated_slide_dots(self) -> None:
        topics = load_topics(SAMPLE_CONFIG)
        slide_states = build_slide_states(topics)

        self.assertEqual([topic.name for topic in topics], ["Intro", "Data", "Modeling", "Outcomes"])
        self.assertEqual(len(slide_states), 14)
        self.assertEqual(slide_states[5], slide_states[6])
        self.assertNotEqual(slide_states[6], slide_states[7])

    def test_adds_grouped_topic_footer_and_styles_active_topic(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            slide_count = add_progress_dots_to_presentation(SAMPLE_DECK, SAMPLE_CONFIG, output_path)
            presentation = Presentation(str(output_path))

            self.assertEqual(slide_count, 14)

            slide_two = presentation.slides[1]
            slide_seven = presentation.slides[6]
            slide_eight = presentation.slides[7]

            self.assertEqual(self.topic_label_color(slide_two, "Intro"), "000000")
            self.assertEqual(self.topic_label_color(slide_two, "Data"), "E7E6E6")
            self.assertEqual(self.topic_dot_line_colors(slide_two, "Intro"), ["000000"] * 5)
            self.assertEqual(self.active_dot_indices(slide_two, "Intro"), [1])
            self.assertEqual(self.topic_dot_line_colors(slide_two, "Data"), ["E7E6E6"] * 2)

            self.assertEqual(self.topic_label_color(slide_seven, "Data"), "000000")
            self.assertEqual(self.topic_label_color(slide_seven, "Intro"), "E7E6E6")
            self.assertEqual(self.active_dot_indices(slide_seven, "Data"), [0])
            self.assertEqual(self.active_dot_indices(slide_eight, "Data"), [1])
            self.assertLess(self.topic_left(slide_two, "Intro"), self.topic_left(slide_two, "Data"))

    def test_rerunning_does_not_duplicate_tracker_shapes(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            first_output = Path(tmpdir) / "first-output.pptx"
            second_output = Path(tmpdir) / "second-output.pptx"

            add_progress_dots_to_presentation(SAMPLE_DECK, SAMPLE_CONFIG, first_output)
            add_progress_dots_to_presentation(first_output, SAMPLE_CONFIG, second_output)

            presentation = Presentation(str(second_output))
            slide = presentation.slides[0]
            labels = [shape for shape in slide.shapes if shape.name.startswith(f"{TRACKER_SHAPE_PREFIX} Topic Label")]
            dots = [shape for shape in slide.shapes if shape.name.startswith(f"{TRACKER_SHAPE_PREFIX} Topic Dot")]

            self.assertEqual(len(labels), 4)
            self.assertEqual(len(dots), 12)

    def topic_label_color(self, slide, topic_name: str) -> str:
        label = self.named_shape(slide, f"{TRACKER_SHAPE_PREFIX} Topic Label {topic_name}")
        for paragraph in label.text_frame.paragraphs:
            for run in paragraph.runs:
                return str(run.font.color.rgb)
        raise AssertionError(f"No text color found for topic label {topic_name}.")

    def topic_dot_line_colors(self, slide, topic_name: str) -> list[str]:
        dots = self.topic_dots(slide, topic_name)
        return [str(dot.line.color.rgb) for dot in dots]

    def active_dot_indices(self, slide, topic_name: str) -> list[int]:
        indices = []
        for index, dot in enumerate(self.topic_dots(slide, topic_name)):
            if str(dot.fill.fore_color.rgb) == "000000":
                indices.append(index)
        return indices

    def topic_left(self, slide, topic_name: str) -> int:
        label = self.named_shape(slide, f"{TRACKER_SHAPE_PREFIX} Topic Label {topic_name}")
        return label.left

    def topic_dots(self, slide, topic_name: str) -> list:
        dots = [
            shape
            for shape in slide.shapes
            if shape.name.startswith(f"{TRACKER_SHAPE_PREFIX} Topic Dot {topic_name} ")
        ]
        dots.sort(key=lambda shape: shape.left)
        return dots

    def named_shape(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        raise AssertionError(f"Shape '{name}' not found.")


if __name__ == "__main__":
    unittest.main()
